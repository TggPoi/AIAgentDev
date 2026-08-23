from collections import defaultdict
from itertools import zip_longest
from pathlib import Path
from typing import Protocol

from openpyxl import load_workbook
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from fast_app.domain.knowledge_models import (
    ExcelFieldValue,
    ExcelRow,
    ExcelSheet,
    LoadedDocument,
    LoadedExcelDocument,
    LoadedPowerPointDocument,
    PowerPointSlide,
    VisionImageContent,
    VisionImageOccurrence,
)
from fast_app.ingestion.processing.metadata_models import build_document_metadata

# 文档读取层 loader

class BaseDocumentLoader(Protocol):
    """文档 Loader 的最小协议：把目录内容转换为领域文档列表。"""

    def load(self, base_dir: str) -> list[LoadedDocument]:
        """读取指定知识库目录。"""

        pass

# 从知识库目录递归读取 .md 文件 保留 source_path 使用 UTF-8 读取中文 Markdown 返回原始 MarkdownDocument
class MarkdownDocumentLoader:
    """读取知识库目录中的 UTF-8 Markdown 文件。"""

    def load(self, base_dir: str) -> list[LoadedDocument]:
        """按路径稳定排序并保留每个 Markdown 文件的 source_path。"""

        root = Path(base_dir)
        documents: list[LoadedDocument] = []

        for path in sorted(root.rglob("*.md")):
            source_path = path.as_posix()
            documents.append(
                LoadedDocument(
                    source_path=source_path,
                    content=path.read_text(encoding="utf-8"),
                    document_type="markdown",
                    metadata=build_document_metadata(
                        source_path=source_path,
                        document_type="markdown",
                        knowledge_base_dir=base_dir,
                    ),
                )
            )

        return documents


class TextDocumentLoader:
    """读取知识库目录中的 UTF-8 纯文本文件。"""

    def load(self, base_dir: str) -> list[LoadedDocument]:
        """按路径稳定排序并保留每个文本文件的 source_path。"""

        root = Path(base_dir)
        documents: list[LoadedDocument] = []

        for path in sorted(root.rglob("*.txt")):
            source_path = path.as_posix()
            documents.append(
                LoadedDocument(
                    source_path=source_path,
                    content=path.read_text(encoding="utf-8"),
                    document_type="text",
                    metadata=build_document_metadata(
                        source_path=source_path,
                        document_type="text",
                        knowledge_base_dir=base_dir,
                    ),
                )
            )

        return documents


class PowerPointDocumentLoader:
    """把 PPTX 的可检索文本转换为单个 Markdown-like 文档。"""

    def __init__(
        self,
        *,
        max_image_bytes: int | None = None,
        max_image_pixels: int | None = None,
    ) -> None:
        self._max_image_bytes = max_image_bytes
        self._max_image_pixels = max_image_pixels

    def load(self, base_dir: str) -> list[LoadedDocument]:
        """按路径稳定排序，递归读取目录中的所有 PPTX 文件。"""

        root = Path(base_dir)
        return [
            self.load_file(
                path,
                source_path=path.as_posix(),
                knowledge_base_dir=base_dir,
            )
            for path in sorted(root.rglob("*.pptx"))
        ]

    def load_file(
        self,
        path: str | Path,
        *,
        source_path: str | None = None,
        knowledge_base_dir: str | None = None,
    ) -> LoadedDocument:
        """兼容旧链路，把结构化 PPT 重新渲染为一个 Markdown-like 文档。"""

        document = self.load_structured_file(
            path,
            source_path=source_path,
            knowledge_base_dir=knowledge_base_dir,
        )
        sections: list[str] = []
        for slide in document.slides:
            sections.append(
                f"# Slide {slide.slide_number}: {slide.title or f'Slide {slide.slide_number}'}"
            )
            if slide.content:
                sections.append(slide.content)
            if slide.notes:
                sections.extend(["## Notes", slide.notes])
        return LoadedDocument(
            source_path=document.source_path,
            content="\n\n".join(sections),
            document_type="powerpoint",
            metadata=document.metadata,
        )

    def load_structured_file(
        self,
        path: str | Path,
        *,
        source_path: str | None = None,
        knowledge_base_dir: str | None = None,
    ) -> LoadedPowerPointDocument:
        """提取稳定 slide_id、页面正文、表格和备注。"""

        file_path = Path(path)
        source = source_path or file_path.as_posix()
        presentation = Presentation(file_path)
        metadata = build_document_metadata(
            source_path=source,
            document_type="powerpoint",
            knowledge_base_dir=knowledge_base_dir,
        )
        doc_id = str(metadata["doc_id"])
        warnings: set[str] = set()
        slides: list[PowerPointSlide] = []
        vision_contents: dict[str, VisionImageContent] = {}
        vision_occurrences: list[VisionImageOccurrence] = []

        for slide_index, slide in enumerate(presentation.slides, start=1):
            title_shape = slide.shapes.title
            if title_shape is None:
                # 企业模板常用普通文本框绘制标题；取视觉顺序中的首个有效文本框。
                title_shape = _find_first_powerpoint_text_shape(slide.shapes)
            title = title_shape.text.strip() if title_shape is not None else ""
            # 标题已经写入 section 标题，后续遍历时跳过，避免同一文本重复入库。
            skip_ids = {title_shape.shape_id} if title_shape is not None else set()
            slide_warnings: set[str] = set()
            slide_occurrences: list[VisionImageOccurrence] = []
            body = _extract_powerpoint_shapes(
                slide.shapes,
                skip_ids,
                slide_warnings,
                doc_id=doc_id,
                slide_id=int(getattr(slide, "slide_id", slide_index)),
                slide_number=slide_index,
                vision_contents=vision_contents,
                vision_occurrences=slide_occurrences,
                max_image_bytes=self._max_image_bytes,
                max_image_pixels=self._max_image_pixels,
            )
            vision_occurrences.extend(slide_occurrences)
            warnings.update(slide_warnings)

            # 先判断 has_notes_slide，避免访问 slide.notes_slide 时隐式创建备注页。
            notes = ""
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame is not None:
                    notes = notes_slide.notes_text_frame.text.strip()
            slides.append(
                PowerPointSlide(
                    # 正常 python-pptx Slide 始终提供 slide_id；回退值仅兼容精简测试桩。
                    slide_id=int(getattr(slide, "slide_id", slide_index)),
                    slide_number=slide_index,
                    title=title,
                    content="\n\n".join(part for part in body if part),
                    notes=notes,
                    warnings=tuple(sorted(slide_warnings)),
                    vision_occurrence_ids=tuple(
                        occurrence.occurrence_id for occurrence in slide_occurrences
                    ),
                )
            )

        metadata.update(
            slide_count=len(presentation.slides),
            image_occurrence_count=len(vision_occurrences),
            extraction_warnings=sorted(warnings),
        )
        return LoadedPowerPointDocument(
            source_path=source,
            slides=slides,
            vision_contents=vision_contents,
            vision_occurrences=vision_occurrences,
            metadata=metadata,
        )


def _extract_powerpoint_shapes(
    shapes,
    skip_shape_ids: set[int],
    warnings: set[str],
    *,
    doc_id: str,
    slide_id: int,
    slide_number: int,
    vision_contents: dict[str, VisionImageContent],
    vision_occurrences: list[VisionImageOccurrence],
    max_image_bytes: int | None,
    max_image_pixels: int | None,
) -> list[str]:
    """按视觉位置遍历一层 Shape Tree，并递归提取组合图形的内容。"""

    result: list[str] = []
    # shape_id 用作同位置图形的稳定兜底键，保证重复解析产生相同 section 顺序。
    ordered_shapes = sorted(
        shapes,
        key=lambda shape: (int(shape.top), int(shape.left), int(shape.shape_id)),
    )
    for shape in ordered_shapes:
        if shape.shape_id in skip_shape_ids:
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # GroupShape 自身也是 Shape Tree；每一层都重新排序后递归处理。
            result.extend(
                _extract_powerpoint_shapes(
                    shape.shapes,
                    skip_shape_ids,
                    warnings,
                    doc_id=doc_id,
                    slide_id=slide_id,
                    slide_number=slide_number,
                    vision_contents=vision_contents,
                    vision_occurrences=vision_occurrences,
                    max_image_bytes=max_image_bytes,
                    max_image_pixels=max_image_pixels,
                )
            )
            continue
        if getattr(shape, "has_table", False):
            result.append(_powerpoint_table_to_markdown(shape.table))
            continue
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text:
                result.append(text)
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                raw = bytes(shape.image.blob)
                content = VisionImageContent.from_raw(
                    raw,
                    media_type=str(
                        getattr(shape.image, "content_type", "")
                        or "application/octet-stream"
                    ),
                    max_bytes=max_image_bytes,
                    max_pixels=max_image_pixels,
                )
                vision_contents.setdefault(content.content_id, content)
                occurrence_index = len(vision_occurrences) + 1
                vision_occurrences.append(
                    VisionImageOccurrence(
                        occurrence_id=(
                            f"imgocc:ppt:{doc_id}:{slide_id}:{shape.shape_id}:"
                            f"{content.content_id}"
                        ),
                        content_id=content.content_id,
                        source_locator=(
                            f"slide[{slide_number}]/shape[{int(shape.shape_id)}]"
                        ),
                        page_or_slide_number=slide_number,
                        anchor_id=str(shape.shape_id),
                        occurrence_index=occurrence_index,
                    )
                )
            except Exception:
                warnings.add("PPT_IMAGE_EXTRACTION_FAILED")
            continue
        if shape.shape_type in {
            MSO_SHAPE_TYPE.CHART,
            MSO_SHAPE_TYPE.DIAGRAM,
            MSO_SHAPE_TYPE.MEDIA,
        }:
            warnings.add("PPT_UNSUPPORTED_VISUAL_SKIPPED")
    return result


def _find_first_powerpoint_text_shape(shapes):
    """按视觉顺序递归查找可作为无占位符页面标题的首个文本 Shape。"""

    for shape in sorted(
        shapes,
        key=lambda item: (int(item.top), int(item.left), int(item.shape_id)),
    ):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            nested = _find_first_powerpoint_text_shape(shape.shapes)
            if nested is not None:
                return nested
            continue
        if getattr(shape, "has_text_frame", False) and shape.text.strip():
            return shape
    return None


def _powerpoint_table_to_markdown(table) -> str:
    """把 PPT 表格转换为 Markdown 表格，首行作为检索友好的表头。"""

    rows = [[_escape_markdown(cell.text) for cell in row.cells] for row in table.rows]
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    normalized = [row + [""] * (width - len(row)) for row in rows]
    header = normalized[0]
    body = normalized[1:]
    lines = [
        f"| {' | '.join(header)} |",
        f"| {' | '.join(['---'] * width)} |",
    ]
    lines.extend(f"| {' | '.join(row)} |" for row in body)
    return "\n".join(lines)


class ExcelDocumentLoader:
    """把 XLSX 的可见工作表转换为保留原始坐标的 Markdown-like 文档。"""

    row_block_size = 100
    max_non_empty_cells = 100_000

    def load(self, base_dir: str) -> list[LoadedDocument]:
        """按路径稳定排序，递归读取目录中的所有 XLSX 文件。"""

        root = Path(base_dir)
        return [
            self.load_file(
                path,
                source_path=path.as_posix(),
                knowledge_base_dir=base_dir,
            )
            for path in sorted(root.rglob("*.xlsx"))
        ]

    def load_file(
        self,
        path: str | Path,
        *,
        source_path: str | None = None,
        knowledge_base_dir: str | None = None,
    ) -> LoadedDocument:
        """兼容旧链路，把结构化工作表重新渲染为 Markdown-like 文档。"""

        document = self.load_structured_file(
            path,
            source_path=source_path,
            knowledge_base_dir=knowledge_base_dir,
        )
        sections: list[str] = []
        for sheet in document.sheets:
            blocks: dict[int, list[tuple[int, dict[str, str]]]] = defaultdict(list)
            for row in sheet.rows:
                block_start = (
                    (row.row_number - 1) // self.row_block_size
                ) * self.row_block_size + 1
                blocks[block_start].append((row.row_number, row.values))
            sections.extend(
                _excel_sheet_to_markdown(
                    sheet.name,
                    blocks,
                    sheet.business_header_hint,
                    self.row_block_size,
                )
            )
        return LoadedDocument(
            source_path=document.source_path,
            content="\n\n".join(sections),
            document_type="spreadsheet",
            metadata=document.metadata,
        )

    def load_structured_file(
        self,
        path: str | Path,
        *,
        source_path: str | None = None,
        knowledge_base_dir: str | None = None,
    ) -> LoadedExcelDocument:
        """用公式视图和缓存值视图配对解析结构化工作表。"""

        file_path = Path(path)
        source = source_path or file_path.as_posix()
        # openpyxl 的一个 Workbook 只能选择公式或缓存值，因此必须独立打开两次。
        formula_wb = load_workbook(
            file_path,
            read_only=True,
            data_only=False,
            keep_links=False,
        )
        value_wb = load_workbook(
            file_path,
            read_only=True,
            data_only=True,
            keep_links=False,
        )
        warnings: set[str] = set()
        header_hints: dict[str, dict[str, str]] = {}
        sheets: list[ExcelSheet] = []
        non_empty_cells = 0
        sheet_names = list(formula_wb.sheetnames)

        try:
            for sheet_name in formula_wb.sheetnames:
                # 两个 Workbook 按相同工作表名、相同行列坐标并行读取。
                formula_ws = formula_wb[sheet_name]
                value_ws = value_wb[sheet_name]
                if formula_ws.sheet_state != "visible":
                    warnings.add("xlsx_hidden_sheet_skipped")
                    continue
                # 部分生产工具不写 worksheet dimension；强制计算一次，避免同一文件
                # 经 Excel 保存后 max_column 从 None 变为实际值而误触发全 Sheet 更新。
                if formula_ws.max_column is None:
                    formula_ws.calculate_dimension(force=True)

                rows: list[ExcelRow] = []
                for formula_row, value_row in zip_longest(
                    formula_ws.iter_rows(),
                    value_ws.iter_rows(),
                    fillvalue=(),
                ):
                    values: dict[str, str] = {}
                    cells: dict[str, ExcelFieldValue] = {}
                    row_number = 0
                    for formula_cell, value_cell in zip_longest(
                        formula_row,
                        value_row,
                        fillvalue=None,
                    ):
                        cell = next(
                            (
                                candidate
                                for candidate in (formula_cell, value_cell)
                                if candidate is not None and hasattr(candidate, "row")
                            ),
                            None,
                        )
                        if cell is None:
                            continue
                        row_number = int(cell.row)
                        # A/B/C 是权威列标识；业务表头只会作为额外提示保存。
                        column = get_column_letter(int(cell.column))
                        rendered, cache_missing = _render_excel_cell(
                            formula_cell,
                            value_cell,
                        )
                        if cache_missing:
                            warnings.add("xlsx_formula_cache_missing")
                        if rendered == "":
                            continue
                        non_empty_cells += 1
                        if non_empty_cells > self.max_non_empty_cells:
                            raise ValueError(
                                "XLSX 非空单元格超过 100000 个，拒绝导入"
                            )
                        values[column] = rendered
                        formula = (
                            str(formula_cell.value)
                            if formula_cell is not None
                            and formula_cell.data_type == "f"
                            else None
                        )
                        cached = (
                            _escape_markdown(value_cell.value)
                            if formula is not None
                            and value_cell is not None
                            and value_cell.value is not None
                            else None
                        )
                        cells[column] = ExcelFieldValue(
                            value=rendered,
                            formula=formula,
                            cached_value=cached,
                            source_column=column,
                            source_coordinate=f"{column}{row_number}",
                        )

                    if not values:
                        continue
                    # 第一条非空行仅作为检索提示，不参与坐标定位或权限判断。
                    if sheet_name not in header_hints:
                        header_hints[sheet_name] = dict(values)
                    rows.append(
                        ExcelRow(row_number=row_number, values=values, cells=cells)
                    )

                sheets.append(
                    ExcelSheet(
                        name=sheet_name,
                        rows=rows,
                        business_header_hint=header_hints.get(sheet_name, {}),
                        source_columns=[
                            get_column_letter(column)
                            for column in range(1, int(formula_ws.max_column or 0) + 1)
                        ],
                    )
                )
        finally:
            # 两个只读 Workbook 各自持有文件资源，异常路径也必须分别关闭。
            formula_wb.close()
            value_wb.close()

        metadata = build_document_metadata(
            source_path=source,
            document_type="spreadsheet",
            knowledge_base_dir=knowledge_base_dir,
        )
        metadata.update(
            sheet_names=sheet_names,
            business_header_hints=header_hints,
            extraction_warnings=sorted(warnings),
        )
        return LoadedExcelDocument(
            source_path=source,
            sheets=sheets,
            metadata=metadata,
        )


def _render_excel_cell(formula_cell, value_cell) -> tuple[str, bool]:
    """渲染单元格；公式同时保留表达式和 Excel 最后保存的缓存值。"""

    formula_value = None if formula_cell is None else formula_cell.value
    cached_value = None if value_cell is None else value_cell.value
    if formula_cell is not None and formula_cell.data_type == "f":
        formula_text = str(formula_value)
        if cached_value is None:
            return formula_text, True
        return f"{formula_text} => {_escape_markdown(cached_value)}", False
    return _escape_markdown(formula_value), False


def _excel_sheet_to_markdown(
    sheet_name: str,
    blocks: dict[int, list[tuple[int, dict[str, str]]]],
    header_hint: dict[str, str],
    block_size: int,
) -> list[str]:
    """把工作表的稀疏行块渲染为带行号、列坐标和表头提示的 Markdown。"""

    if not blocks:
        return []
    result = [f"# Sheet: {sheet_name}"]
    hint = ", ".join(f"{column}={value}" for column, value in header_hint.items())
    if hint:
        result.append(f"Business header hint: {hint}")
    for block_start in sorted(blocks):
        rows = blocks[block_start]
        columns = sorted(
            {column for _, values in rows for column in values},
            key=lambda column: _excel_column_number(column),
        )
        result.append(f"## Rows {block_start}-{block_start + block_size - 1}")
        result.append(f"| Row | {' | '.join(columns)} |")
        result.append(f"| --- | {' | '.join(['---'] * len(columns))} |")
        for row_number, values in rows:
            result.append(
                f"| {row_number} | {' | '.join(values.get(column, '') for column in columns)} |"
            )
    return result


def _excel_column_number(column: str) -> int:
    """把 Excel 列字母转换为整数，用于 A、Z、AA 的自然顺序排序。"""

    result = 0
    for char in column:
        result = result * 26 + ord(char) - ord("A") + 1
    return result


def _escape_markdown(value: object) -> str:
    """转义表格分隔符并压平换行，防止单元格破坏 Markdown 表格结构。"""

    if value is None:
        return ""
    return str(value).replace("|", "\\|").replace("\r\n", " ").replace("\n", " ").strip()

# 组合 loader
class CompositeDocumentLoader:
    """顺序调用多个 Loader，并把结果合并为一个文档列表。"""

    def __init__(self, loaders: list[BaseDocumentLoader]):
        """保存需要依次执行的 Loader。"""

        self.loaders = loaders

    def load(self, base_dir: str) -> list[LoadedDocument]:
        """用同一知识库目录调用所有 Loader 并合并结果。"""

        documents: list[LoadedDocument] = []

        for loader in self.loaders:
            documents.extend(loader.load(base_dir))

        return documents


def build_default_document_loader() -> CompositeDocumentLoader:
    """构建 Markdown 主链路使用的 Markdown/TXT Loader。"""

    return CompositeDocumentLoader(
        loaders=[
            MarkdownDocumentLoader(),
            TextDocumentLoader(),
        ]
    )


