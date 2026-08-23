"""验证知识文档页面读取与共享 Document Access Policy 使用相同 ACL。"""

from __future__ import annotations

import asyncio
import hashlib
import io
import json
from datetime import UTC, datetime
from unittest.mock import AsyncMock

from docx import Document
from fastapi import FastAPI
from fastapi.testclient import TestClient
from openpyxl import Workbook
from pptx import Presentation
from pypdf import PdfWriter
from sqlalchemy import text

from fast_app.api.knowledge_document_routes import router
from fast_app.core.config import get_settings
from fast_app.core.exception_handlers import register_exception_handlers
from fast_app.db.gitlab_tables import GitLabDocumentTable, GitLabSourceTable
from fast_app.db.session import create_database_engine, create_session_factory
from fast_app.dependencies.knowledge_document_dependencies import (
    get_knowledge_document_read_service,
)
from fast_app.dependencies.user_context import get_current_user_context
from fast_app.domain.auth_models import AccountType
from fast_app.domain.user_context import CurrentUserContext
from fast_app.schemas.knowledge_document_schema import (
    KnowledgeDocumentContentResponse,
    KnowledgeDocumentDetail,
    KnowledgeDocumentItem,
    KnowledgeDocumentListResponse,
)
from fast_app.services.auth.auth_crypto import hash_password
from fast_app.services.exceptions import (
    AuthenticationError,
    KnowledgeDocumentContentTooLargeError,
    KnowledgeDocumentContentUnavailableError,
    KnowledgeDocumentNotFoundError,
)
from fast_app.services.knowledge.document_access_policy import DocumentAccessPolicy
from fast_app.services.knowledge.document_access_repository import (
    DocumentAccessRepository,
)
from fast_app.services.knowledge.knowledge_document_read_repository import (
    KnowledgeDocumentReadRepository,
    KnowledgeDocumentRecord,
)
from fast_app.services.knowledge.knowledge_document_read_service import (
    KnowledgeDocumentDownload,
    KnowledgeDocumentReadService,
    _extract_text_preview,
)


ADMIN_ID = "user_knowledge_read_admin"
ART_USER_ID = "user_knowledge_read_art"
SOURCE_ART_ID = "source_knowledge_read_art"
SOURCE_DEV_ID = "source_knowledge_read_dev"
DOC_ART = "doc_knowledge_read_art"
DOC_GRANTED = "doc_knowledge_read_granted"
DOC_HIDDEN = "doc_knowledge_read_hidden"
DOC_PUBLIC = "doc_knowledge_read_public"
DOC_ORIGINAL_ACL = "doc_knowledge_read_original_acl"
REVISION = "a" * 40
RAW_CONTENT = {
    DOC_ART: b"# Art Guide\n\nOwn department.",
    DOC_GRANTED: b"# Granted Guide\n\nExact cross-department grant.",
    DOC_HIDDEN: b"# Hidden Guide\n\nMust not leak.",
    DOC_PUBLIC: b"Public text document.",
    DOC_ORIGINAL_ACL: b"# Legacy ACL\n\nExplicit original ACL.",
}


class FakeContentGateway:
    def __init__(self, content: dict[str, bytes]) -> None:
        self.content = content
        self.calls: list[tuple[str, str, str]] = []

    async def fetch(self, *, source, document) -> bytes | None:
        self.calls.append(
            (document.doc_id, document.repository_path, document.source_revision)
        )
        return self.content.get(document.doc_id)


def main() -> None:
    asyncio.run(assert_database_flow())
    assert_preview_formats()
    assert_http_contract()
    print("knowledge_document_read=passed")


async def assert_database_flow() -> None:
    settings = get_settings()
    engine = create_database_engine(settings)
    session_factory = create_session_factory(engine)
    try:
        async with session_factory() as session:
            await _cleanup(session)
            await _seed_facts(session)
            gateway = FakeContentGateway(dict(RAW_CONTENT))
            service = KnowledgeDocumentReadService(
                settings=settings,
                repository=KnowledgeDocumentReadRepository(session),
                access_policy=DocumentAccessPolicy(
                    DocumentAccessRepository(session)
                ),
                content_gateway=gateway,  # type: ignore[arg-type]
            )
            art_user = _art_user()

            first_page = await service.list_documents(
                art_user,
                cursor=None,
                limit=2,
                query=None,
                department_code=None,
                document_type=None,
            )
            assert len(first_page.items) == 2
            assert first_page.next_cursor is not None
            second_page = await service.list_documents(
                art_user,
                cursor=first_page.next_cursor,
                limit=10,
                query=None,
                department_code=None,
                document_type=None,
            )
            items = [*first_page.items, *second_page.items]
            assert {item.doc_id for item in items} == {
                DOC_ART,
                DOC_GRANTED,
                DOC_PUBLIC,
                DOC_ORIGINAL_ACL,
            }
            assert len({item.doc_id for item in items}) == 4
            assert DOC_HIDDEN not in {item.doc_id for item in items}
            assert {item.doc_id: item.access_source for item in items} == {
                DOC_ART: "department",
                DOC_GRANTED: "explicit_grant",
                DOC_PUBLIC: "public",
                DOC_ORIGINAL_ACL: "original_acl",
            }
            admin_page = await service.list_documents(
                CurrentUserContext(
                    user_id=ADMIN_ID,
                    username="knowledge_read_admin",
                    account_type=AccountType.ADMIN,
                    is_authenticated=True,
                    auth_source="jwt",
                    global_role_codes=["system_admin"],
                ),
                cursor=None,
                limit=20,
                query=None,
                department_code=None,
                document_type=None,
            )
            assert {item.doc_id for item in admin_page.items} == set(RAW_CONTENT)
            assert all(item.access_source == "admin" for item in admin_page.items)

            development = await service.list_documents(
                art_user,
                cursor=None,
                limit=20,
                query=None,
                department_code="development",
                document_type=None,
            )
            assert {item.doc_id for item in development.items} == {
                DOC_GRANTED,
                DOC_PUBLIC,
                DOC_ORIGINAL_ACL,
            }
            query_result = await service.list_documents(
                art_user,
                cursor=None,
                limit=20,
                query="granted",
                department_code=None,
                document_type="markdown",
            )
            assert [item.doc_id for item in query_result.items] == [DOC_GRANTED]

            detail = await service.get_detail(art_user, DOC_GRANTED)
            assert detail.access_source == "explicit_grant"
            assert detail.source_revision == REVISION
            assert detail.department_code == "development"
            try:
                await service.get_detail(art_user, DOC_HIDDEN)
            except KnowledgeDocumentNotFoundError:
                pass
            else:
                raise AssertionError("已知 doc_id 绕过了列表 ACL")

            content = await service.get_content(art_user, DOC_GRANTED)
            assert content.render_mode == "markdown"
            assert content.source_revision == REVISION
            assert "Exact cross-department grant" in content.content
            download = await service.get_download(art_user, DOC_GRANTED)
            assert download.content == RAW_CONTENT[DOC_GRANTED]
            assert download.file_name == "granted.md"
            assert download.media_type == "text/markdown"
            assert gateway.calls[-1] == (
                DOC_GRANTED,
                "development/granted.md",
                REVISION,
            )

            try:
                await service.list_documents(
                    CurrentUserContext(
                        user_id="anonymous",
                        is_authenticated=False,
                        auth_source="anonymous",
                    ),
                    cursor=None,
                    limit=20,
                    query=None,
                    department_code=None,
                    document_type=None,
                )
            except AuthenticationError:
                pass
            else:
                raise AssertionError("匿名用户读取了知识文档目录")

            await session.execute(
                text(
                    "update document_access_grants "
                    "set status = 'revoked', revoked_by_user_id = :actor_id, "
                    "revoked_at = now() where grantee_user_id = :user_id "
                    "and doc_id = :doc_id and status = 'active'"
                ),
                {
                    "actor_id": ADMIN_ID,
                    "user_id": ART_USER_ID,
                    "doc_id": DOC_GRANTED,
                },
            )
            await session.commit()
            after_revoke = await service.list_documents(
                art_user,
                cursor=None,
                limit=20,
                query=None,
                department_code=None,
                document_type=None,
            )
            assert DOC_GRANTED not in {item.doc_id for item in after_revoke.items}
            try:
                await service.get_download(art_user, DOC_GRANTED)
            except KnowledgeDocumentNotFoundError:
                pass
            else:
                raise AssertionError("撤销后仍可通过下载接口读取文档")

            too_large_gateway = FakeContentGateway(
                {DOC_ART: b"x" * (settings.gitlab_source_file_max_bytes + 1)}
            )
            too_large_service = KnowledgeDocumentReadService(
                settings=settings,
                repository=KnowledgeDocumentReadRepository(session),
                access_policy=DocumentAccessPolicy(
                    DocumentAccessRepository(session)
                ),
                content_gateway=too_large_gateway,  # type: ignore[arg-type]
            )
            try:
                await too_large_service.get_download(art_user, DOC_ART)
            except KnowledgeDocumentContentTooLargeError:
                pass
            else:
                raise AssertionError("超限源文件未被拒绝")

            mismatch_service = KnowledgeDocumentReadService(
                settings=settings,
                repository=KnowledgeDocumentReadRepository(session),
                access_policy=DocumentAccessPolicy(
                    DocumentAccessRepository(session)
                ),
                content_gateway=FakeContentGateway(
                    {DOC_ART: b"same-size-content-mismatch"}
                ),  # type: ignore[arg-type]
            )
            try:
                await mismatch_service.get_download(art_user, DOC_ART)
            except KnowledgeDocumentContentUnavailableError:
                pass
            else:
                raise AssertionError("manifest blob 与下载内容不一致时未被拒绝")
    finally:
        async with session_factory() as cleanup_session:
            await _cleanup(cleanup_session)
        await engine.dispose()


async def _seed_facts(session) -> None:
    await session.execute(
        text(
            """
            insert into users (id, username, password_hash, status)
            values
                (:admin_id, 'knowledge_read_admin', :password_hash, 'active'),
                (:art_user_id, 'knowledge_read_art', :password_hash, 'active')
            """
        ),
        {
            "admin_id": ADMIN_ID,
            "art_user_id": ART_USER_ID,
            "password_hash": hash_password("KnowledgeRead123!"),
        },
    )
    await session.execute(
        text(
            """
            insert into user_departments (id, user_id, department_code, is_primary)
            values ('user_dept_knowledge_read_art', :user_id, 'art', true)
            """
        ),
        {"user_id": ART_USER_ID},
    )
    for source_id, project_id, project_path, department in (
        (SOURCE_ART_ID, 92001, "knowledge/art", "art"),
        (SOURCE_DEV_ID, 92002, "knowledge/development", "development"),
    ):
        await session.execute(
            text(
                """
                insert into gitlab_sources
                    (id, base_url, host_id, project_id, project_path,
                     target_branch, department_code, default_visibility,
                     sync_token_env, agent_token_env, webhook_secret_env, status)
                values
                    (:id, 'https://gitlab.example.test', 'knowledge-read-test',
                     :project_id, :project_path, 'main', :department, 'department',
                     'SYNC_TOKEN', 'AGENT_TOKEN', 'WEBHOOK_SECRET', 'active')
                """
            ),
            {
                "id": source_id,
                "project_id": project_id,
                "project_path": project_path,
                "department": department,
            },
        )
    documents = (
        (DOC_ART, SOURCE_ART_ID, "art/guide.md", "markdown", "department", []),
        (
            DOC_GRANTED,
            SOURCE_DEV_ID,
            "development/granted.md",
            "markdown",
            "department",
            [],
        ),
        (
            DOC_HIDDEN,
            SOURCE_DEV_ID,
            "development/hidden.md",
            "markdown",
            "department",
            [],
        ),
        (
            DOC_PUBLIC,
            SOURCE_DEV_ID,
            "development/public.txt",
            "text",
            "public",
            [],
        ),
        (
            DOC_ORIGINAL_ACL,
            SOURCE_DEV_ID,
            "development/original-acl.md",
            "markdown",
            "restricted",
            [ART_USER_ID],
        ),
    )
    for doc_id, source_id, path, doc_type, visibility, allowed_users in documents:
        raw = RAW_CONTENT[doc_id]
        department = "art" if source_id == SOURCE_ART_ID else "development"
        acl = {
            "visibility": visibility,
            "allowed_departments": [department] if visibility != "public" else [],
            "allowed_users": allowed_users,
        }
        await session.execute(
            text(
                """
                insert into gitlab_documents
                    (doc_id, source_id, repository_path, blob_id, source_revision,
                     content_hash, acl_hash, parser_version,
                     chunk_strategy_version, chunk_config_fingerprint,
                     document_type, acl_json, status)
                values
                    (:doc_id, :source_id, :path, :blob_id, :revision,
                     :content_hash, :acl_hash, 'parser-v1', 'chunk-v1',
                     'config-v1', :document_type, cast(:acl_json as jsonb), 'active')
                """
            ),
            {
                "doc_id": doc_id,
                "source_id": source_id,
                "path": path,
                "blob_id": _git_blob_id(raw),
                "revision": REVISION,
                "content_hash": hashlib.sha256(raw).hexdigest(),
                "acl_hash": hashlib.sha256(
                    json.dumps(acl, sort_keys=True).encode("utf-8")
                ).hexdigest(),
                "document_type": doc_type,
                "acl_json": json.dumps(acl),
            },
        )
    await session.execute(
        text(
            """
            insert into document_access_grants
                (id, grantee_user_id, doc_id, granted_by_user_id, status)
            values
                ('grant_knowledge_read', :user_id, :doc_id, :admin_id, 'active')
            """
        ),
        {
            "user_id": ART_USER_ID,
            "doc_id": DOC_GRANTED,
            "admin_id": ADMIN_ID,
        },
    )
    await session.commit()


async def _cleanup(session) -> None:
    await session.execute(
        text(
            "delete from document_access_grants "
            "where grantee_user_id = :user_id or granted_by_user_id = :admin_id"
        ),
        {"user_id": ART_USER_ID, "admin_id": ADMIN_ID},
    )
    await session.execute(
        text("delete from gitlab_documents where doc_id = any(:doc_ids)"),
        {"doc_ids": list(RAW_CONTENT)},
    )
    await session.execute(
        text("delete from gitlab_sources where id = any(:source_ids)"),
        {"source_ids": [SOURCE_ART_ID, SOURCE_DEV_ID]},
    )
    await session.execute(
        text("delete from users where id = any(:user_ids)"),
        {"user_ids": [ADMIN_ID, ART_USER_ID]},
    )
    await session.commit()


def assert_preview_formats() -> None:
    settings = get_settings()
    samples: list[tuple[str, bytes, str]] = []

    docx_buffer = io.BytesIO()
    docx = Document()
    docx.add_heading("Word Preview", level=1)
    docx.add_paragraph("DOCX body text")
    docx.save(docx_buffer)
    samples.append(("word", docx_buffer.getvalue(), "DOCX body text"))

    pptx_buffer = io.BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "PowerPoint Preview"
    slide.placeholders[1].text = "PPTX body text"
    presentation.save(pptx_buffer)
    samples.append(("powerpoint", pptx_buffer.getvalue(), "PPTX body text"))

    xlsx_buffer = io.BytesIO()
    workbook = Workbook()
    worksheet = workbook.active
    worksheet["A1"] = "Spreadsheet Preview"
    worksheet["A2"] = "XLSX body text"
    workbook.save(xlsx_buffer)
    workbook.close()
    samples.append(("spreadsheet", xlsx_buffer.getvalue(), "XLSX body text"))

    pdf_buffer = io.BytesIO()
    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    writer.write(pdf_buffer)
    samples.append(("pdf", pdf_buffer.getvalue(), ""))

    for document_type, raw, expected in samples:
        record = _preview_record(document_type)
        mode, content, warnings = _extract_text_preview(
            record,
            raw,
            settings,
        )
        assert mode == "extracted_text"
        assert expected in content
        if document_type == "pdf":
            assert "pdf_text_content_unavailable" in warnings


def _preview_record(document_type: str) -> KnowledgeDocumentRecord:
    suffix = {
        "word": "docx",
        "powerpoint": "pptx",
        "spreadsheet": "xlsx",
        "pdf": "pdf",
    }[document_type]
    return KnowledgeDocumentRecord(
        document=GitLabDocumentTable(
            doc_id=f"preview-{document_type}",
            source_id="preview-source",
            repository_path=f"preview/source.{suffix}",
            source_revision=REVISION,
            content_hash="content",
            acl_hash="acl",
            parser_version="parser",
            chunk_strategy_version="chunk",
            chunk_config_fingerprint="config",
            document_type=document_type,
            acl_json={"visibility": "public", "allowed_users": []},
            status="active",
        ),
        source=GitLabSourceTable(
            id="preview-source",
            base_url="https://gitlab.example.test",
            host_id="preview",
            project_id=1,
            project_path="preview/project",
            target_branch="main",
            department_code="preview",
            default_visibility="public",
            sync_token_env="SYNC_TOKEN",
            agent_token_env="AGENT_TOKEN",
            webhook_secret_env="WEBHOOK_SECRET",
            status="active",
        ),
    )


def assert_http_contract() -> None:
    now = datetime.now(UTC)
    item = KnowledgeDocumentItem(
        doc_id="doc-http",
        title="安全 文档",
        file_name="安全 文档.md",
        repository_path="development/安全 文档.md",
        department_code="development",
        document_type="markdown",
        source_revision=REVISION,
        updated_at=now,
        access_source="department",
    )
    fake_service = AsyncMock()
    fake_service.list_documents.return_value = KnowledgeDocumentListResponse(
        items=[item]
    )
    fake_service.get_detail.return_value = KnowledgeDocumentDetail(
        **item.model_dump(),
        source_id="source-http",
        source_project_path="knowledge/development",
        visibility="department",
    )
    fake_service.get_content.return_value = KnowledgeDocumentContentResponse(
        doc_id="doc-http",
        source_revision=REVISION,
        document_type="markdown",
        render_mode="markdown",
        content="# 安全文档",
        truncated=False,
        warnings=[],
    )
    fake_service.get_download.return_value = KnowledgeDocumentDownload(
        content=b"# safe",
        file_name='安全 "文档".md',
        media_type="text/markdown",
        source_revision=REVISION,
    )
    app = FastAPI()
    register_exception_handlers(app)
    app.include_router(router)
    app.dependency_overrides[get_current_user_context] = _art_user
    app.dependency_overrides[get_knowledge_document_read_service] = (
        lambda: fake_service
    )
    with TestClient(app) as client:
        assert client.get("/knowledge/documents").status_code == 200
        assert client.get("/knowledge/documents/doc-http").status_code == 200
        assert (
            client.get("/knowledge/documents/doc-http/content").status_code
            == 200
        )
        response = client.get("/knowledge/documents/doc-http/download")
        assert response.status_code == 200
        assert response.content == b"# safe"
        assert response.headers["x-content-type-options"] == "nosniff"
        assert "filename*=UTF-8''" in response.headers["content-disposition"]
        assert "\r" not in response.headers["content-disposition"]
        assert "\n" not in response.headers["content-disposition"]
        assert client.get("/knowledge/documents/" + "x" * 65).status_code == 422
    schema = app.openapi()
    for path in (
        "/knowledge/documents",
        "/knowledge/documents/{doc_id}",
        "/knowledge/documents/{doc_id}/content",
        "/knowledge/documents/{doc_id}/download",
    ):
        assert path in schema["paths"]


def _art_user() -> CurrentUserContext:
    return CurrentUserContext(
        user_id=ART_USER_ID,
        username="knowledge_read_art",
        account_type=AccountType.EMPLOYEE,
        is_authenticated=True,
        auth_source="jwt",
        department_codes=["art"],
        primary_department_code="art",
    )


def _git_blob_id(content: bytes) -> str:
    return hashlib.sha1(
        f"blob {len(content)}\0".encode("utf-8") + content
    ).hexdigest()


if __name__ == "__main__":
    main()
