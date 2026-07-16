from __future__ import annotations

import asyncio
import hashlib
import sys
from contextlib import suppress
from datetime import UTC, datetime
from pathlib import Path
from tempfile import TemporaryDirectory
from types import SimpleNamespace
from unittest.mock import patch
from zipfile import ZIP_DEFLATED, ZipFile

from fastapi import FastAPI
from fastapi.testclient import TestClient
from elasticsearch import AsyncElasticsearch
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from pymilvus import MilvusClient
from xlsxwriter import Workbook as XlsxWriterWorkbook
from sqlalchemy import delete, select, update

ROOT_DIR = Path(__file__).resolve().parents[1]
SRC_DIR = ROOT_DIR / "src"
if str(SRC_DIR) not in sys.path:
    sys.path.insert(0, str(SRC_DIR))

from fast_app.core.config import Settings, get_settings
from fast_app.components.embeddings.mock_embedding_client import MockEmbeddingClient
from fast_app.core.exception_handlers import register_exception_handlers
from fast_app.api.knowledge_import_routes import (
    ExcelProfileConfirmRequest,
    _validate_excel_profile_request,
    get_import_job_repository,
    router,
)
from fast_app.dependencies.rag_dependencies import get_permission_service
from fast_app.dependencies.user_context import get_current_user_context
from fast_app.domain.agent_tool_permissions import (
    DepartmentPermissionScope,
    EffectivePermissionSet,
    PermissionCode,
)
from fast_app.domain.user_context import CurrentUserContext
from fast_app.domain.knowledge_models import (
    ExcelRow,
    ExcelSheet,
    LoadedExcelDocument,
    LoadedPowerPointDocument,
    PowerPointSlide,
)
from fast_app.db.ingestion_tables import KnowledgeDocumentTable, KnowledgeIngestionJobTable
from fast_app.db.auth_tables import UserTable
from fast_app.db.session import create_database_engine, create_session_factory
from fast_app.ingestion.chunk_builders import ChunkBuildOptions, MarkdownChunkBuilder
from fast_app.ingestion.document_loaders import (
    ExcelDocumentLoader,
    MarkdownDocumentLoader,
    PowerPointDocumentLoader,
    TextDocumentLoader,
    build_default_document_loader,
)
from fast_app.ingestion.import_jobs import (
    HEARTBEAT_SECONDS,
    LEASE_SECONDS,
    ImportJobValidationError,
    KnowledgeImportJobRepository,
    new_import_job_id,
    normalize_upload_filename,
)
from fast_app.ingestion.incremental_store import (
    ExistingChunkState,
    apply_chunk_diff,
    build_chunk_diff,
    load_es_chunk_states,
    load_milvus_chunk_states,
    verify_chunk_convergence,
)
from fast_app.ingestion.metadata_models import build_doc_id, build_document_metadata
from fast_app.ingestion.ooxml_validation import OOXMLValidationError, validate_ooxml_package
from fast_app.ingestion.office_chunk_builders import (
    ExcelChunkBuilder,
    ExcelConfigurationRequired,
    PowerPointChunkBuilder,
)
from fast_app.ingestion.rag_store_writer import replace_docs_rag_stores
from fast_app.ingestion.worker import KnowledgeImportWorker
from fast_app.middlewares.request_size_middleware import RequestSizeLimitMiddleware


def build_pptx(path: Path, image_path: Path) -> None:
    """生成包含正文、两级组合图形、表格、图片和备注的测试 PPTX。"""

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Product Plan"
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1)).text = "Body"
    group = slide.shapes.add_group_shape()
    group.shapes.add_textbox(Inches(2), Inches(2), Inches(2), Inches(1)).text = "Grouped"
    nested = group.shapes.add_group_shape()
    nested.shapes.add_textbox(Inches(3), Inches(3), Inches(2), Inches(1)).text = "Nested"
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(4), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "Asset"
    table.cell(0, 1).text = "Owner"
    table.cell(1, 0).text = "Server"
    table.cell(1, 1).text = "Ops"
    slide.shapes.add_picture(str(image_path), Inches(6), Inches(1), Inches(1), Inches(1))
    slide.notes_slide.notes_text_frame.text = "Speaker notes"
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(path)


def build_incremental_pptx(path: Path, *, insert_middle: bool = False) -> None:
    """生成可保持旧 slide_id 的两页 PPT，并可在中间插入第三页。"""

    presentation = Presentation()
    for title, body in (("One", "alpha"), ("Two", "beta")):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = title
        slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(3), Inches(1)
        ).text = body
    if insert_middle:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Inserted"
        slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(3), Inches(1)
        ).text = "new page"
        # 只调整 sldId 顺序，旧页面的 OOXML slide_id 不会改变。
        slide_ids = presentation.slides._sldIdLst
        inserted = slide_ids[-1]
        slide_ids.remove(inserted)
        slide_ids.insert(1, inserted)
    presentation.save(path)


def build_xlsx(path: Path) -> None:
    """生成包含表头、缓存公式、跨区段行和隐藏表的测试 XLSX。"""

    workbook = XlsxWriterWorkbook(path)
    sheet = workbook.add_worksheet("Assets")
    sheet.write_row(0, 0, ["ID", "Name", "Cost|CNY"])
    sheet.write_row(1, 0, [1, "Server", 10])
    sheet.write_formula(1, 3, "=A2+1", None, 2)
    sheet.write(100, 0, 101)
    hidden = workbook.add_worksheet("Hidden")
    hidden.hide()
    hidden.write(0, 0, "secret")
    workbook.close()


def test_office_builders_and_incremental_diff() -> None:
    """覆盖 PPT 插页、Excel 行列移动和 metadata-only 向量复用。"""

    options = ChunkBuildOptions(
        source="test",
        max_chars=4000,
        overlap_chars=0,
        max_tokens=2000,
        min_chars=1,
    )
    fingerprint = "embedding-v1"
    ppt_metadata = build_document_metadata(
        source_path="knowledge/hr/plan.pptx", document_type="powerpoint"
    )
    original_ppt = LoadedPowerPointDocument(
        source_path="knowledge/hr/plan.pptx",
        metadata=ppt_metadata,
        slides=[
            PowerPointSlide(256, 1, "A", "alpha"),
            PowerPointSlide(257, 2, "B", "beta"),
        ],
    )
    inserted_ppt = LoadedPowerPointDocument(
        source_path=original_ppt.source_path,
        metadata=dict(ppt_metadata),
        slides=[
            PowerPointSlide(256, 1, "A", "alpha"),
            PowerPointSlide(999, 2, "NEW", "new"),
            PowerPointSlide(257, 3, "B", "beta"),
        ],
    )
    ppt_builder = PowerPointChunkBuilder()
    before = ppt_builder.build(
        original_ppt, options, embedding_fingerprint=fingerprint
    )
    after = ppt_builder.build(inserted_ppt, options, embedding_fingerprint=fingerprint)
    before_by_slide = {chunk.metadata["slide_id"]: chunk for chunk in before}
    after_by_slide = {chunk.metadata["slide_id"]: chunk for chunk in after}
    assert before_by_slide[256].id == after_by_slide[256].id
    assert before_by_slide[257].id == after_by_slide[257].id
    assert before_by_slide[257].metadata["content_hash"] == after_by_slide[257].metadata["content_hash"]
    assert before_by_slide[257].metadata["index_hash"] != after_by_slide[257].metadata["index_hash"]

    excel_metadata = build_document_metadata(
        source_path="knowledge/hr/assets.xlsx", document_type="spreadsheet"
    )
    profile = {
        "mode": "record",
        "sheets": [
            {
                "sheet_key": "assets",
                "sheet_name_aliases": ["Assets"],
                "header_row": 1,
                "identity_field_ids": ["asset_id"],
                "fields": [
                    {
                        "field_id": "asset_id",
                        "display_name": "ID",
                        "header_aliases": ["资产编号"],
                        "required": True,
                        "indexed": True,
                    },
                    {
                        "field_id": "name",
                        "display_name": "Name",
                        "header_aliases": ["名称"],
                        "required": True,
                        "indexed": True,
                    },
                ],
            }
        ],
    }
    original_excel = LoadedExcelDocument(
        source_path="knowledge/hr/assets.xlsx",
        metadata=excel_metadata,
        sheets=[
            ExcelSheet(
                "Assets",
                [
                    ExcelRow(1, {"A": "ID", "B": "Name"}),
                    ExcelRow(2, {"A": "1", "B": "Server"}),
                    ExcelRow(3, {"A": "2", "B": "Switch"}),
                ],
                {"A": "ID", "B": "Name"},
            )
        ],
    )
    # 插入行并移动物理列；字段序列仍由 Profile 决定。
    moved_excel = LoadedExcelDocument(
        source_path=original_excel.source_path,
        metadata=dict(excel_metadata),
        sheets=[
            ExcelSheet(
                "Assets",
                [
                    ExcelRow(1, {"A": "Name", "C": "ID"}),
                    ExcelRow(3, {"A": "Server", "C": "1"}),
                    ExcelRow(4, {"A": "Switch", "C": "2"}),
                ],
                {"A": "Name", "C": "ID"},
            )
        ],
    )
    excel_builder = ExcelChunkBuilder()
    excel_before = excel_builder.build(
        original_excel, options, profile=profile, embedding_fingerprint=fingerprint
    )
    excel_after = excel_builder.build(
        moved_excel, options, profile=profile, embedding_fingerprint=fingerprint
    )
    before_by_identity = {
        chunk.metadata["row_identity"]: chunk for chunk in excel_before
    }
    after_by_identity = {chunk.metadata["row_identity"]: chunk for chunk in excel_after}
    assert set(before_by_identity) == set(after_by_identity)
    for identity in before_by_identity:
        assert before_by_identity[identity].id == after_by_identity[identity].id
        assert (
            before_by_identity[identity].metadata["content_hash"]
            == after_by_identity[identity].metadata["content_hash"]
        )
        assert (
            before_by_identity[identity].metadata["index_hash"]
            != after_by_identity[identity].metadata["index_hash"]
        )

    unknown_column = LoadedExcelDocument(
        source_path=original_excel.source_path,
        metadata=dict(excel_metadata),
        sheets=[
            ExcelSheet(
                "Assets",
                [
                    ExcelRow(1, {"A": "ID", "B": "Name", "C": "Owner"}),
                    ExcelRow(2, {"A": "1", "B": "Server", "C": "Ops"}),
                ],
                {"A": "ID", "B": "Name", "C": "Owner"},
            )
        ],
    )
    try:
        excel_builder.build(
            unknown_column,
            options,
            profile=profile,
            embedding_fingerprint=fingerprint,
        )
    except ExcelConfigurationRequired:
        pass
    else:
        raise AssertionError("有值未知列必须进入 awaiting_configuration")

    section_before_document = LoadedExcelDocument(
        source_path=original_excel.source_path,
        metadata=dict(excel_metadata),
        sheets=[
            ExcelSheet(
                "Assets",
                [ExcelRow(1, {"A": "ID"}), ExcelRow(101, {"A": "101"})],
                {"A": "ID"},
                ["A"],
            ),
            ExcelSheet(
                "Owners", [ExcelRow(1, {"A": "Ops"})], {"A": "Ops"}, ["A"]
            ),
        ],
    )
    section_after_document = LoadedExcelDocument(
        source_path=original_excel.source_path,
        metadata=dict(excel_metadata),
        sheets=[
            ExcelSheet(
                "Assets",
                [ExcelRow(1, {"A": "ID"}), ExcelRow(101, {"A": "101"})],
                {"A": "ID"},
                ["A", "B"],
            ),
            ExcelSheet(
                "Owners", [ExcelRow(1, {"A": "Ops"})], {"A": "Ops"}, ["A"]
            ),
        ],
    )
    section_before = excel_builder.build(
        section_before_document,
        options,
        profile={"mode": "section"},
        embedding_fingerprint=fingerprint,
    )
    section_after = excel_builder.build(
        section_after_document,
        options,
        profile={"mode": "section"},
        embedding_fingerprint=fingerprint,
    )
    before_sections = {chunk.id: chunk for chunk in section_before}
    after_sections = {chunk.id: chunk for chunk in section_after}
    assert set(before_sections) == set(after_sections)
    assert next(
        chunk.metadata["source_columns"]
        for chunk in section_after
        if chunk.metadata["sheet_name"] == "Assets"
    ) == ["A", "B"]
    for chunk_id, chunk in before_sections.items():
        changed = chunk.metadata["sheet_name"] == "Assets"
        assert (
            chunk.metadata["index_hash"] != after_sections[chunk_id].metadata["index_hash"]
        ) is changed

    unchanged = excel_before[0]
    metadata_only = excel_after[1]
    es_states = {
        unchanged.id: ExistingChunkState(dict(unchanged.metadata)),
        metadata_only.id: ExistingChunkState(dict(excel_before[1].metadata)),
    }
    milvus_states = {
        unchanged.id: ExistingChunkState(dict(unchanged.metadata), [0.1, 0.2]),
        metadata_only.id: ExistingChunkState(
            dict(excel_before[1].metadata), [0.3, 0.4]
        ),
    }
    diff = build_chunk_diff(
        [unchanged, metadata_only], es_states, milvus_states, embedding_dim=2
    )
    assert diff.counts["unchanged"] == 1
    assert diff.counts["metadata_only"] == 1
    assert diff.counts["embedded"] == 0

    # Section 模式在可容纳单行时不得从表格行中间切开。
    row_document = LoadedExcelDocument(
        source_path=original_excel.source_path,
        metadata=dict(excel_metadata),
        sheets=[
            ExcelSheet(
                "Rows",
                [
                    ExcelRow(1, {"A": "ID", "B": "Name"}),
                    ExcelRow(2, {"A": "1", "B": "A" * 24}),
                    ExcelRow(3, {"A": "2", "B": "B" * 24}),
                ],
                {"A": "ID", "B": "Name"},
                ["A", "B"],
            )
        ],
    )
    row_chunks = excel_builder.build(
        row_document,
        ChunkBuildOptions("test", 120, 0, 120, 1),
        profile={"mode": "section"},
        embedding_fingerprint=fingerprint,
    )
    data_lines = [
        line
        for chunk in row_chunks
        for line in chunk.content.splitlines()
        if line.startswith(("| 2 |", "| 3 |"))
    ]
    assert data_lines == [
        "| 2 | 1 | " + "A" * 24 + " |",
        "| 3 | 2 | " + "B" * 24 + " |",
    ]
    assert all(
        line.startswith(("Sheet: ", "| "))
        for chunk in row_chunks
        for line in chunk.content.splitlines()
    )
    for chunk in row_chunks:
        data_rows = [
            int(line.split("|")[1].strip())
            for line in chunk.content.splitlines()
            if line.startswith("| ") and not line.startswith("| ---")
            and line.split("|")[1].strip().isdigit()
        ]
        assert chunk.metadata["row_start"] == min(data_rows)
        assert chunk.metadata["row_end"] == max(data_rows)
        assert chunk.metadata["source_columns"] == ["A", "B"]


def test_mixed_excel_profile_and_coordinates() -> None:
    """验证同一 Workbook 混合模式、精确坐标和旧 Profile 兼容。"""

    options = ChunkBuildOptions("test", 4000, 0, 2000, 1)
    fingerprint = "embedding-v1"
    metadata = build_document_metadata(
        source_path="knowledge/development/assets.xlsx",
        document_type="spreadsheet",
    )
    original = LoadedExcelDocument(
        source_path="knowledge/development/assets.xlsx",
        metadata=metadata,
        sheets=[
            ExcelSheet(
                "资产清单",
                [
                    ExcelRow(1, {"A": "资产ID", "K": "负责人", "N": "优先级"}),
                    ExcelRow(3, {"A": "AST-0002", "K": "赵凯", "N": "P1"}),
                ],
                {"A": "资产ID", "K": "负责人", "N": "优先级"},
                ["A", "K", "N"],
            ),
            ExcelSheet(
                "分类说明",
                [
                    ExcelRow(2, {"A": "角色资产", "B": "标准 Skeleton"}),
                    ExcelRow(3, {"A": "动画资产", "B": "Anim Montage"}),
                ],
                {"A": "角色资产", "B": "标准 Skeleton"},
                ["A", "B"],
            ),
        ],
    )
    profile = {
        "mode": "mixed",
        "sheets": [
            {
                "sheet_key": "asset_list",
                "mode": "record",
                "sheet_name_aliases": ["资产清单"],
                "header_row": 1,
                "identity_field_ids": ["asset_id"],
                "fields": [
                    {
                        "field_id": "asset_id",
                        "display_name": "资产ID",
                        "required": True,
                        "indexed": True,
                    },
                    {
                        "field_id": "owner",
                        "display_name": "负责人",
                        "required": True,
                        "indexed": True,
                    },
                    {
                        "field_id": "priority",
                        "display_name": "优先级",
                        "required": True,
                        "indexed": True,
                    },
                ],
            },
            {
                "sheet_key": "category_guide",
                "mode": "section",
                "sheet_name_aliases": ["分类说明"],
            },
        ],
    }
    builder = ExcelChunkBuilder()
    before = builder.build(
        original, options, profile=profile, embedding_fingerprint=fingerprint
    )
    record = next(chunk for chunk in before if chunk.metadata["excel_mode"] == "record")
    section = next(chunk for chunk in before if chunk.metadata["excel_mode"] == "section")
    assert record.metadata["row_identity"] == "AST-0002"
    assert record.metadata["row_number"] == 3
    assert record.metadata["field_coordinates"] == {
        "asset_id": "A3",
        "owner": "K3",
        "priority": "N3",
    }
    assert section.metadata["sheet_key"] == "category_guide"
    assert section.metadata["row_start"] == 2
    assert section.metadata["row_end"] == 3
    assert section.metadata["source_columns"] == ["A", "B"]

    moved = LoadedExcelDocument(
        source_path=original.source_path,
        metadata=dict(metadata),
        sheets=[
            ExcelSheet(
                "资产清单",
                [
                    ExcelRow(1, {"A": "负责人", "C": "资产ID", "D": "优先级"}),
                    ExcelRow(4, {"A": "赵凯", "C": "AST-0002", "D": "P1"}),
                ],
                {"A": "负责人", "C": "资产ID", "D": "优先级"},
                ["A", "C", "D"],
            ),
            original.sheets[1],
        ],
    )
    after = builder.build(
        moved, options, profile=profile, embedding_fingerprint=fingerprint
    )
    moved_record = next(
        chunk for chunk in after if chunk.metadata["excel_mode"] == "record"
    )
    assert moved_record.id == record.id
    assert moved_record.metadata["content_hash"] == record.metadata["content_hash"]
    assert moved_record.metadata["index_hash"] != record.metadata["index_hash"]
    assert moved_record.metadata["row_number"] == 4
    assert moved_record.metadata["field_coordinates"] == {
        "asset_id": "C4",
        "owner": "A4",
        "priority": "D4",
    }
    es_states = {chunk.id: ExistingChunkState(dict(chunk.metadata)) for chunk in before}
    milvus_states = {
        chunk.id: ExistingChunkState(dict(chunk.metadata), [0.1, 0.2])
        for chunk in before
    }
    diff = build_chunk_diff(after, es_states, milvus_states, embedding_dim=2)
    assert diff.counts["metadata_only"] == 1
    assert diff.counts["embedded"] == 0

    changed = LoadedExcelDocument(
        source_path=moved.source_path,
        metadata=dict(metadata),
        sheets=[
            ExcelSheet(
                "资产清单",
                [
                    ExcelRow(1, {"A": "负责人", "C": "资产ID", "D": "优先级"}),
                    ExcelRow(4, {"A": "赵凯", "C": "AST-0002", "D": "P0"}),
                ],
                moved.sheets[0].business_header_hint,
                ["A", "C", "D"],
            ),
            original.sheets[1],
        ],
    )
    changed_record = next(
        chunk
        for chunk in builder.build(
            changed, options, profile=profile, embedding_fingerprint=fingerprint
        )
        if chunk.metadata["excel_mode"] == "record"
    )
    assert changed_record.id == record.id
    assert changed_record.metadata["content_hash"] != record.metadata["content_hash"]

    missing_sheet_profile = {"mode": "mixed", "sheets": profile["sheets"][:1]}
    try:
        builder.build(
            original,
            options,
            profile=missing_sheet_profile,
            embedding_fingerprint=fingerprint,
        )
    except ExcelConfigurationRequired:
        pass
    else:
        raise AssertionError("Mixed Profile 不得静默跳过未配置 Sheet")

    # API schema 保留旧模式，并拒绝 mixed 配置缺失 Sheet mode。
    _validate_excel_profile_request(
        ExcelProfileConfirmRequest.model_validate(
            {
                "preview_fingerprint": "a" * 64,
                "mode": "mixed",
                "profile_name": "Mixed assets",
                "sheets": profile["sheets"],
            }
        )
    )
    _validate_excel_profile_request(
        ExcelProfileConfirmRequest.model_validate(
            {
                "preview_fingerprint": "a" * 64,
                "mode": "section",
                "profile_name": "Legacy section",
                "sheets": [],
            }
        )
    )
    invalid = dict(profile["sheets"][0])
    invalid.pop("mode")
    try:
        _validate_excel_profile_request(
            ExcelProfileConfirmRequest.model_validate(
                {
                    "preview_fingerprint": "a" * 64,
                    "mode": "mixed",
                    "profile_name": "Invalid mixed",
                    "sheets": [invalid],
                }
            )
        )
    except ImportJobValidationError:
        pass
    else:
        raise AssertionError("Mixed Sheet 缺少 mode 时必须拒绝")

    duplicate_key = dict(profile["sheets"][1])
    duplicate_key["sheet_key"] = profile["sheets"][0]["sheet_key"]
    try:
        _validate_excel_profile_request(
            ExcelProfileConfirmRequest.model_validate(
                {
                    "preview_fingerprint": "a" * 64,
                    "mode": "mixed",
                    "profile_name": "Duplicate sheet key",
                    "sheets": [profile["sheets"][0], duplicate_key],
                }
            )
        )
    except ImportJobValidationError:
        pass
    else:
        raise AssertionError("Mixed Profile 的 sheet_key 重复时必须拒绝")

    invalid_identity = dict(profile["sheets"][0])
    invalid_identity["identity_field_ids"] = ["missing_field"]
    try:
        _validate_excel_profile_request(
            ExcelProfileConfirmRequest.model_validate(
                {
                    "preview_fingerprint": "a" * 64,
                    "mode": "mixed",
                    "profile_name": "Invalid identity",
                    "sheets": [invalid_identity, profile["sheets"][1]],
                }
            )
        )
    except ImportJobValidationError:
        pass
    else:
        raise AssertionError("Record Sheet 引用不存在的主键字段时必须拒绝")


def test_default_loader_excludes_office() -> None:
    """默认 Markdown ingestion 只接收 Markdown/TXT，Office 交给异步 Worker。"""

    loaders = build_default_document_loader().loaders
    assert [type(loader) for loader in loaders] == [
        MarkdownDocumentLoader,
        TextDocumentLoader,
    ]


def test_office_loaders_and_stable_chunks(root: Path) -> None:
    """验证 Office 文本提取、警告、坐标保留和稳定 Chunk ID。"""

    image_path = root / "image.png"
    Image.new("RGB", (2, 2), color="red").save(image_path)
    pptx_path = root / "plan.pptx"
    xlsx_path = root / "assets.xlsx"
    build_pptx(pptx_path, image_path)
    build_xlsx(xlsx_path)

    validate_ooxml_package(pptx_path)
    validate_ooxml_package(xlsx_path)

    pptx_document = PowerPointDocumentLoader().load_file(pptx_path)
    assert all(
        text in pptx_document.content
        for text in ("Product Plan", "Body", "Grouped", "Nested", "Server", "Speaker notes")
    )
    assert "pptx_visual_content_skipped" in pptx_document.metadata["extraction_warnings"]

    no_placeholder = Presentation()
    slide = no_placeholder.slides.add_slide(no_placeholder.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(1)).text = "Body below"
    slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(3), Inches(1)).text = "Inferred title"
    inferred_path = root / "inferred-title.pptx"
    no_placeholder.save(inferred_path)
    inferred = PowerPointDocumentLoader().load_structured_file(inferred_path)
    assert inferred.slides[0].title == "Inferred title"
    assert "Inferred title" not in inferred.slides[0].content
    assert "Body below" in inferred.slides[0].content

    xlsx_document = ExcelDocumentLoader().load_file(xlsx_path)
    assert "| Row | A | B | C | D |" in xlsx_document.content
    assert "=A2+1 => 2" in xlsx_document.content
    assert "Rows 101-200" in xlsx_document.content
    assert "Cost\\|CNY" in xlsx_document.content
    assert xlsx_document.metadata["business_header_hints"]["Assets"]["A"] == "ID"
    structured_xlsx = ExcelDocumentLoader().load_structured_file(xlsx_path)
    assert structured_xlsx.sheets[0].source_columns == ["A", "B", "C", "D"]
    assert "xlsx_hidden_sheet_skipped" in xlsx_document.metadata["extraction_warnings"]
    assert "xlsx_formula_cache_missing" not in xlsx_document.metadata["extraction_warnings"]

    formula_only = root / "formula-only.xlsx"
    workbook = Workbook()
    workbook.active["A1"] = "=1+1"
    workbook.save(formula_only)
    formula_document = ExcelDocumentLoader().load_file(formula_only)
    assert "xlsx_formula_cache_missing" in formula_document.metadata["extraction_warnings"]

    options = ChunkBuildOptions(
        source="test",
        max_chars=500,
        overlap_chars=20,
        max_tokens=500,
        min_chars=1,
    )
    builder = MarkdownChunkBuilder()
    first = builder.build([xlsx_document], options)
    second = builder.build([ExcelDocumentLoader().load_file(xlsx_path)], options)
    assert [chunk.id for chunk in first] == [chunk.id for chunk in second]
    assert any("Sheet: Assets" in chunk.metadata["section_path"] for chunk in first)


def test_notes_none_guard(root: Path) -> None:
    """验证备注页存在但 notes_text_frame 缺失时仍能安全解析。"""

    class Shapes(list):
        title = None

    class NotesSlide:
        notes_text_frame = None

    class Slide:
        shapes = Shapes()
        has_notes_slide = True
        notes_slide = NotesSlide()

    class FakePresentation:
        slides = [Slide()]

    with patch("fast_app.ingestion.document_loaders.Presentation", return_value=FakePresentation()):
        document = PowerPointDocumentLoader().load_file(root / "unused.pptx")
    assert "Slide 1" in document.content


def test_ooxml_rejections(root: Path) -> None:
    """验证损坏 ZIP、核心文件缺失和目录穿越使用稳定错误码拒绝。"""

    corrupt = root / "corrupt.xlsx"
    corrupt.write_bytes(b"PK\x03\x04broken")
    try:
        validate_ooxml_package(corrupt)
        raise AssertionError("corrupt ZIP should fail")
    except OOXMLValidationError as exc:
        assert exc.code == "INVALID_OOXML_PACKAGE"

    missing_core = root / "missing.xlsx"
    with ZipFile(missing_core, "w", ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", "<Types/>")
        archive.writestr("_rels/.rels", "<Relationships/>")
    try:
        validate_ooxml_package(missing_core)
        raise AssertionError("missing xl/workbook.xml should fail")
    except OOXMLValidationError as exc:
        assert exc.code == "OOXML_CORE_FILE_MISSING"

    unsafe = root / "unsafe.xlsx"
    with ZipFile(unsafe, "w", ZIP_DEFLATED) as archive:
        archive.writestr("../escape.xml", "x")
    try:
        validate_ooxml_package(unsafe)
        raise AssertionError("unsafe ZIP path should fail")
    except OOXMLValidationError as exc:
        assert exc.code == "OOXML_UNSAFE_PATH"


def test_filename_and_request_limits() -> None:
    """验证安全文件名以及普通请求与上传请求使用不同全局上限。"""

    assert normalize_upload_filename("资产列表.xlsx") == ("资产列表.xlsx", "spreadsheet")
    for filename in ("../bad.xlsx", "CON.xlsx", ".hidden.pptx", "bad.pdf"):
        try:
            normalize_upload_filename(filename)
            raise AssertionError(f"unsafe filename accepted: {filename}")
        except ImportJobValidationError:
            pass

    settings = Settings(_env_file=None)
    assert settings.max_request_body_bytes == 64 * 1024
    assert settings.max_upload_file_bytes == 20 * 1024 * 1024
    assert settings.max_upload_request_body_bytes == 21 * 1024 * 1024

    app = FastAPI()
    app.add_middleware(
        RequestSizeLimitMiddleware,
        max_body_bytes=64,
        max_upload_body_bytes=100,
    )

    @app.post("/normal")
    async def normal() -> dict[str, bool]:
        return {"ok": True}

    @app.post("/knowledge-documents/import-jobs")
    async def upload() -> dict[str, bool]:
        return {"ok": True}

    @app.post("/knowledge-documents/doc_0123456789abcdef/import-jobs")
    async def update_upload() -> dict[str, bool]:
        return {"ok": True}

    client = TestClient(app)
    assert client.post("/normal", content=b"x" * 65).status_code == 413
    assert client.post("/knowledge-documents/import-jobs", content=b"x" * 65).status_code == 200
    assert client.post("/knowledge-documents/import-jobs/", content=b"x" * 101).status_code == 413
    assert (
        client.post(
            "/knowledge-documents/doc_0123456789abcdef/import-jobs",
            content=b"x" * 65,
        ).status_code
        == 200
    )
    assert (
        client.post(
            "/knowledge-documents/not-a-doc-id/import-jobs", content=b"x" * 65
        ).status_code
        == 413
    )


def test_api_contract(root: Path) -> None:
    """验证创建、查询、列表 API 契约和任务所有者可见性。"""

    xlsx_path = root / "api-assets.xlsx"
    build_xlsx(xlsx_path)
    row: KnowledgeIngestionJobTable | None = None
    staged_paths: list[Path] = []
    created_values: list[dict] = []

    class FakeRepository:
        async def find_active_target(self, target_path: str):
            return None

        async def create(self, **values):
            nonlocal row
            created_values.append(values)
            staged_paths.append(Path(values["staged_path"]))
            now = datetime.now(UTC)
            row = KnowledgeIngestionJobTable(
                id=values["job_id"],
                user_id=values["user_id"],
                department_code=values["department_code"],
                original_filename=values["original_filename"],
                target_path=values["target_path"],
                staged_path=values["staged_path"],
                sha256=values["sha256"],
                document_type=values["document_type"],
                file_size=values["file_size"],
                status="pending",
                phase="queued",
                attempt_count=0,
                max_attempts=3,
                document_count=0,
                chunk_count=0,
                warnings_json=[],
                request_id=values["request_id"],
                trace_id=values["trace_id"],
                created_at=now,
                updated_at=now,
            )
            return row

        async def get_document(self, doc_id: str):
            documents = {
                "doc_excel": SimpleNamespace(
                    doc_id="doc_excel",
                    source_path=(
                        root / "knowledge-base" / "development" / "assets.xlsx"
                    ).as_posix(),
                    department_code="development",
                    document_type="spreadsheet",
                    current_sha256="b" * 64,
                    status="active",
                ),
                "doc_ppt": SimpleNamespace(
                    doc_id="doc_ppt",
                    source_path=(
                        root / "knowledge-base" / "development" / "plan.pptx"
                    ).as_posix(),
                    department_code="development",
                    document_type="powerpoint",
                    current_sha256="c" * 64,
                    status="active",
                ),
            }
            return documents.get(doc_id)

        async def get_active_excel_profile(self, doc_id: str):
            raise AssertionError("reconfigure=true 不得读取 active Profile")

        async def get(self, job_id: str):
            return row if row is not None and row.id == job_id else None

        async def list_for_user(self, *, user_id, status, limit):
            if row is None or (user_id is not None and row.user_id != user_id):
                return []
            return [row]

    class FakePermissionService:
        async def get_effective_permissions(self, user_id: str):
            return EffectivePermissionSet(
                user_id=user_id,
                department_scopes=[
                    DepartmentPermissionScope(
                        department_code="development",
                        permission_codes={
                            PermissionCode.KNOWLEDGE_DOCUMENT_CREATE,
                            PermissionCode.KNOWLEDGE_DOCUMENT_UPDATE,
                        },
                    )
                ],
            )

    current_user = CurrentUserContext(
        user_id="user_test",
        is_authenticated=True,
        auth_source="jwt",
        role="user",
    )
    repository = FakeRepository()
    settings = Settings(
        _env_file=None,
        KNOWLEDGE_BASE_DIR=str(root / "knowledge-base"),
    )
    app = FastAPI()
    app.include_router(router)
    register_exception_handlers(app)
    app.dependency_overrides[get_import_job_repository] = lambda: repository
    app.dependency_overrides[get_permission_service] = lambda: FakePermissionService()
    app.dependency_overrides[get_current_user_context] = lambda: current_user
    from fast_app.core.config import get_settings

    app.dependency_overrides[get_settings] = lambda: settings

    client = TestClient(app)
    response = client.post(
        "/knowledge-documents/import-jobs",
        data={"department_code": "development"},
        files={
            "file": (
                "assets.xlsx",
                xlsx_path.read_bytes(),
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )
        },
    )
    assert response.status_code == 202, response.text
    payload = response.json()
    assert payload["status"] == "pending"
    assert payload["phase"] == "queued"
    assert payload["target_path"].endswith("development/assets.xlsx")
    job_id = payload["job_id"]
    assert client.get(f"/knowledge-documents/import-jobs/{job_id}").status_code == 200
    assert client.get("/knowledge-documents/import-jobs").json()["items"][0]["job_id"] == job_id

    reconfigure = client.post(
        "/knowledge-documents/doc_excel/import-jobs",
        data={
            "expected_sha256": "b" * 64,
            "reconfigure_excel_profile": "true",
        },
        files={
            "file": (
                "assets.xlsx",
                xlsx_path.read_bytes(),
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )
        },
    )
    assert reconfigure.status_code == 202, reconfigure.text
    assert created_values[-1]["operation"] == "update"
    assert created_values[-1]["excel_profile_snapshot"] is None

    ppt_reconfigure = client.post(
        "/knowledge-documents/doc_ppt/import-jobs",
        data={
            "expected_sha256": "c" * 64,
            "reconfigure_excel_profile": "true",
        },
        files={"file": ("plan.pptx", b"unused", "application/octet-stream")},
    )
    assert ppt_reconfigure.status_code == 422, ppt_reconfigure.text
    assert "KNOWLEDGE_IMPORT_INVALID_FILE" in ppt_reconfigure.text

    app.dependency_overrides[get_current_user_context] = lambda: current_user.model_copy(
        update={"user_id": "other_user"}
    )
    assert client.get(f"/knowledge-documents/import-jobs/{job_id}").status_code == 404
    for staged_path in staged_paths:
        staged_path.unlink(missing_ok=True)


async def test_replace_docs_idempotency(root: Path) -> None:
    """连续替换两次同一文档，确认 ES/Milvus 最终只保留稳定主键集合。"""

    xlsx_path = root / "idempotent.xlsx"
    build_xlsx(xlsx_path)
    document = ExcelDocumentLoader().load_file(
        xlsx_path,
        source_path="knowledge-base/development/idempotent.xlsx",
    )
    chunks = MarkdownChunkBuilder().build(
        [document],
        ChunkBuildOptions(
            source="test",
            max_chars=500,
            overlap_chars=20,
            max_tokens=500,
            min_chars=1,
        ),
    )
    settings = Settings(_env_file=None, EMBEDDING_DIM=3)
    vectors = [[0.0, 0.0, 0.0] for _ in chunks]
    es_store: dict[str, object] = {}
    milvus_store: dict[str, object] = {}

    async def fake_replace_es(*, client, settings, chunks):
        doc_ids = {chunk.metadata["doc_id"] for chunk in chunks}
        es_store_copy = {
            chunk_id: chunk
            for chunk_id, chunk in es_store.items()
            if chunk.metadata["doc_id"] not in doc_ids
        }
        es_store.clear()
        es_store.update(es_store_copy)
        es_store.update({chunk.id: chunk for chunk in chunks})
        return len(chunks), {"deleted": 0}

    def fake_replace_milvus(*, client, settings, chunks, vectors):
        doc_ids = {chunk.metadata["doc_id"] for chunk in chunks}
        milvus_copy = {
            chunk_id: chunk
            for chunk_id, chunk in milvus_store.items()
            if chunk.metadata["doc_id"] not in doc_ids
        }
        milvus_store.clear()
        milvus_store.update(milvus_copy)
        milvus_store.update({chunk.id: chunk for chunk in chunks})
        return {"upsert_count": len(chunks)}, {"delete_count": 0}

    with (
        patch("fast_app.ingestion.rag_store_writer.replace_docs_es_index", fake_replace_es),
        patch(
            "fast_app.ingestion.rag_store_writer.replace_docs_milvus_collection",
            fake_replace_milvus,
        ),
    ):
        await replace_docs_rag_stores(None, None, settings, chunks, vectors)
        await replace_docs_rag_stores(None, None, settings, chunks, vectors)

    expected_ids = {chunk.id for chunk in chunks}
    assert set(es_store) == expected_ids
    assert set(milvus_store) == expected_ids


async def test_real_incremental_stores() -> None:
    """在隔离的真实 ES/Milvus 资源中验证增量写入、删除与单边修复。"""

    suffix = new_import_job_id().replace("import_", "")[:12]
    settings = Settings(
        _env_file=None,
        EMBEDDING_DIM=3,
        ELASTICSEARCH_INDEX_NAME=f"office-incremental-test-{suffix}",
        MILVUS_COLLECTION_NAME=f"office_incremental_test_{suffix}",
    )
    es = AsyncElasticsearch(hosts=[settings.elasticsearch_url])
    milvus = MilvusClient(uri=f"http://{settings.milvus_host}:{settings.milvus_port}")
    options = ChunkBuildOptions(
        source="real-test",
        max_chars=4000,
        overlap_chars=0,
        max_tokens=2000,
        min_chars=1,
    )
    metadata = build_document_metadata(
        source_path="knowledge/real/incremental.pptx", document_type="powerpoint"
    )
    builder = PowerPointChunkBuilder()
    before = builder.build(
        LoadedPowerPointDocument(
            source_path="knowledge/real/incremental.pptx",
            metadata=metadata,
            slides=[
                PowerPointSlide(300, 1, "One", "alpha"),
                PowerPointSlide(301, 2, "Two", "beta"),
            ],
        ),
        options,
        embedding_fingerprint="real-test-v1",
    )
    after = builder.build(
        LoadedPowerPointDocument(
            source_path="knowledge/real/incremental.pptx",
            metadata=dict(metadata),
            slides=[
                PowerPointSlide(300, 1, "One", "alpha"),
                PowerPointSlide(999, 2, "Inserted", "new"),
                PowerPointSlide(301, 3, "Two", "beta"),
            ],
        ),
        options,
        embedding_fingerprint="real-test-v1",
    )
    try:
        first = build_chunk_diff(before, {}, {}, embedding_dim=3)
        await apply_chunk_diff(
            elasticsearch_client=es,
            milvus_client=milvus,
            settings=settings,
            diff=first,
            embedded_vectors=[[0.1, 0.2, 0.3] for _ in first.embed],
        )
        await verify_chunk_convergence(
            elasticsearch_client=es,
            milvus_client=milvus,
            settings=settings,
            chunks=before,
        )
        doc_id = str(metadata["doc_id"])
        es_states = await load_es_chunk_states(es, settings, doc_id)
        milvus_states = load_milvus_chunk_states(milvus, settings, doc_id)
        second = build_chunk_diff(after, es_states, milvus_states, embedding_dim=3)
        assert second.counts["added"] == 1
        assert second.counts["metadata_only"] == 1
        assert second.counts["embedded"] == 1
        await apply_chunk_diff(
            elasticsearch_client=es,
            milvus_client=milvus,
            settings=settings,
            diff=second,
            embedded_vectors=[[0.4, 0.5, 0.6]],
        )
        await verify_chunk_convergence(
            elasticsearch_client=es,
            milvus_client=milvus,
            settings=settings,
            chunks=after,
        )

        # 制造 ES 单边缺失，下一次 diff 必须修复而不是追加重复主键。
        await es.delete(
            index=settings.elasticsearch_index_name,
            id=after[0].id,
            refresh=True,
        )
        repaired = build_chunk_diff(
            after,
            await load_es_chunk_states(es, settings, doc_id),
            load_milvus_chunk_states(milvus, settings, doc_id),
            embedding_dim=3,
        )
        assert repaired.counts["repaired"] == 1
        assert repaired.counts["embedded"] == 0
        await apply_chunk_diff(
            elasticsearch_client=es,
            milvus_client=milvus,
            settings=settings,
            diff=repaired,
            embedded_vectors=[],
        )
        await verify_chunk_convergence(
            elasticsearch_client=es,
            milvus_client=milvus,
            settings=settings,
            chunks=after,
        )
    finally:
        if await es.indices.exists(index=settings.elasticsearch_index_name):
            await es.indices.delete(index=settings.elasticsearch_index_name)
        if milvus.has_collection(settings.milvus_collection_name):
            milvus.drop_collection(settings.milvus_collection_name)
        await es.close()
        milvus.close()


async def test_real_worker_incremental_update() -> None:
    """用真实 PostgreSQL/ES/Milvus 验证 Worker 创建后插页更新只嵌入新页。"""

    suffix = new_import_job_id().replace("import_", "")[:12]
    with TemporaryDirectory() as directory:
        root = Path(directory)
        knowledge_root = root / "knowledge"
        staged_create = root / "create.pptx"
        staged_update = root / "update.pptx"
        build_incremental_pptx(staged_create)
        build_incremental_pptx(staged_update, insert_middle=True)
        update_sha = _test_sha256(staged_update)
        target = knowledge_root / "product_planning" / "plan.pptx"
        target_path = target.as_posix()
        doc_id = build_doc_id(target_path)
        settings = Settings(
            _env_file=None,
            KNOWLEDGE_BASE_DIR=str(knowledge_root),
            EMBEDDING_DIM=3,
            LANGSMITH_TRACING=False,
            ELASTICSEARCH_INDEX_NAME=f"office-worker-test-{suffix}",
            MILVUS_COLLECTION_NAME=f"office_worker_test_{suffix}",
        )
        engine = create_database_engine(settings)
        session_factory = create_session_factory(engine)
        es = AsyncElasticsearch(hosts=[settings.elasticsearch_url])
        milvus = MilvusClient(
            uri=f"http://{settings.milvus_host}:{settings.milvus_port}"
        )

        class CountingEmbeddingClient(MockEmbeddingClient):
            """记录每次 Worker 实际提交给 Embedding 的文本数量。"""

            def __init__(self, dim: int) -> None:
                super().__init__(dim)
                self.batch_sizes: list[int] = []

            async def embed_documents(self, texts: list[str]) -> list[list[float]]:
                self.batch_sizes.append(len(texts))
                return await super().embed_documents(texts)

        embeddings = CountingEmbeddingClient(3)
        worker = KnowledgeImportWorker(
            settings=settings,
            session_factory=session_factory,
            embedding_client=embeddings,
            elasticsearch_client=es,
            milvus_client=milvus,
            worker_id=f"worker-{suffix}",
        )
        create_job_id = new_import_job_id()
        update_job_id = new_import_job_id()
        try:
            async with session_factory() as session:
                user_id = await session.scalar(
                    select(UserTable.id).order_by(UserTable.created_at).limit(1)
                )
                assert user_id is not None
                await KnowledgeImportJobRepository(session).create(
                    job_id=create_job_id,
                    user_id=user_id,
                    department_code="product_planning",
                    original_filename="plan.pptx",
                    target_path=target_path,
                    staged_path=str(staged_create),
                    sha256=_test_sha256(staged_create),
                    document_type="powerpoint",
                    file_size=staged_create.stat().st_size,
                    request_id=create_job_id,
                    trace_id=create_job_id,
                    doc_id=doc_id,
                )
            assert await worker.run_once()
            async with session_factory() as session:
                repository = KnowledgeImportJobRepository(session)
                created = await repository.get(create_job_id)
                document = await repository.get_document(doc_id)
                assert created is not None and created.status == "succeeded"
                assert document is not None and document.version == 1
                base_sha = str(document.current_sha256)

                await repository.create(
                    job_id=update_job_id,
                    user_id=user_id,
                    department_code="product_planning",
                    original_filename="plan.pptx",
                    target_path=target_path,
                    staged_path=str(staged_update),
                    sha256=update_sha,
                    document_type="powerpoint",
                    file_size=staged_update.stat().st_size,
                    request_id=update_job_id,
                    trace_id=update_job_id,
                    doc_id=doc_id,
                    operation="update",
                    base_sha256=base_sha,
                )
            embeddings.batch_sizes.clear()
            assert await worker.run_once()
            async with session_factory() as session:
                repository = KnowledgeImportJobRepository(session)
                updated = await repository.get(update_job_id)
                document = await repository.get_document(doc_id)
                assert updated is not None and updated.status == "succeeded"
                assert document is not None and document.version == 2
                assert updated.diff_counts_json["added"] == 1
                assert updated.diff_counts_json["metadata_only"] == 1
                assert updated.diff_counts_json["embedded"] == 1
                assert embeddings.batch_sizes == [1]
                assert _test_sha256(target) == update_sha
        finally:
            async with session_factory() as session:
                await session.execute(
                    delete(KnowledgeIngestionJobTable).where(
                        KnowledgeIngestionJobTable.id.in_(
                            [create_job_id, update_job_id]
                        )
                    )
                )
                await session.execute(
                    delete(KnowledgeDocumentTable).where(
                        KnowledgeDocumentTable.doc_id == doc_id
                    )
                )
                await session.commit()
            if await es.indices.exists(index=settings.elasticsearch_index_name):
                await es.indices.delete(index=settings.elasticsearch_index_name)
            if milvus.has_collection(settings.milvus_collection_name):
                milvus.drop_collection(settings.milvus_collection_name)
            await es.close()
            milvus.close()
            await engine.dispose()


def _test_sha256(path: Path) -> str:
    """计算测试文件 SHA-256。"""

    return hashlib.sha256(path.read_bytes()).hexdigest()


async def test_periodic_heartbeat() -> None:
    """缩短心跳间隔，验证 Worker 会定期续租并在所有权丢失后停止。"""

    renewals = 0

    class SessionContext:
        async def __aenter__(self):
            return object()

        async def __aexit__(self, exc_type, exc, traceback):
            return False

    class SessionFactory:
        def __call__(self):
            return SessionContext()

    class FakeRepository:
        def __init__(self, session):
            pass

        async def renew_lease(self, job_id: str, worker_id: str) -> bool:
            nonlocal renewals
            renewals += 1
            return renewals < 2

    worker = KnowledgeImportWorker(
        settings=Settings(_env_file=None),
        session_factory=SessionFactory(),
        embedding_client=None,
        elasticsearch_client=None,
        milvus_client=None,
        worker_id="worker-test",
    )
    lease_lost = asyncio.Event()
    with (
        patch("fast_app.ingestion.worker.HEARTBEAT_SECONDS", 0.01),
        patch("fast_app.ingestion.worker.KnowledgeImportJobRepository", FakeRepository),
    ):
        task = asyncio.create_task(worker._heartbeat("job-test", lease_lost))
        await asyncio.wait_for(lease_lost.wait(), timeout=0.2)
        task.cancel()
        with suppress(asyncio.CancelledError):
            await task
    assert renewals == 2
    assert HEARTBEAT_SECONDS == 60
    assert LEASE_SECONDS == 300


def test_worker_department_acl(root: Path) -> None:
    """验证 Loader 结果会被覆盖为任务记录中的服务端可信部门 ACL。"""

    xlsx_path = root / "acl.xlsx"
    build_xlsx(xlsx_path)
    job = KnowledgeIngestionJobTable(
        id="job-acl",
        department_code="product_planning",
        document_type="spreadsheet",
        staged_path=str(xlsx_path),
        target_path="knowledge-base/product_planning/acl.xlsx",
    )
    worker = KnowledgeImportWorker(
        settings=Settings(_env_file=None),
        session_factory=None,
        embedding_client=None,
        elasticsearch_client=None,
        milvus_client=None,
        worker_id="worker-test",
    )
    document = worker._load_document(job)
    assert document.metadata["visibility"] == "department"
    assert document.metadata["allowed_departments"] == ["product_planning"]
    assert document.metadata["permission_source"] == "import_job_department"


async def test_postgres_claim_and_recovery() -> None:
    """用真实 PostgreSQL 验证并发领取、过期租约回收和旧所有者失效。"""

    settings = get_settings()
    engine = create_database_engine(settings)
    session_factory = create_session_factory(engine)
    job_id = new_import_job_id()
    try:
        async with session_factory() as session:
            user_id = await session.scalar(
                select(UserTable.id).order_by(UserTable.created_at).limit(1)
            )
            assert user_id is not None
            target_path = f"knowledge-base/product_planning/{job_id}.xlsx"
            await KnowledgeImportJobRepository(session).create(
                job_id=job_id,
                user_id=user_id,
                department_code="product_planning",
                original_filename=f"{job_id}.xlsx",
                target_path=target_path,
                staged_path=f"runtime/knowledge-imports/{job_id}.xlsx",
                sha256="0" * 64,
                document_type="spreadsheet",
                file_size=1,
                request_id=job_id,
                trace_id=job_id,
                doc_id=build_doc_id(target_path),
            )
        async with session_factory() as session:
            await session.execute(
                update(KnowledgeIngestionJobTable)
                .where(KnowledgeIngestionJobTable.id == job_id)
                .values(created_at=datetime(2000, 1, 1, tzinfo=UTC))
            )
            await session.commit()

        async def claim(worker_id: str):
            async with session_factory() as session:
                return await KnowledgeImportJobRepository(session).claim_next(worker_id)

        first, second = await asyncio.gather(claim("worker-a"), claim("worker-b"))
        claimed = first or second
        assert claimed is not None
        assert (first is None) != (second is None)
        old_worker = claimed.row.worker_id
        assert old_worker is not None

        async with session_factory() as session:
            await session.execute(
                update(KnowledgeIngestionJobTable)
                .where(KnowledgeIngestionJobTable.id == job_id)
                .values(lease_expires_at=datetime(2000, 1, 1, tzinfo=UTC))
            )
            await session.commit()

        recovered = await claim("worker-recovery")
        assert recovered is not None
        assert recovered.row.attempt_count == 2
        async with session_factory() as session:
            repository = KnowledgeImportJobRepository(session)
            assert not await repository.renew_lease(job_id, old_worker)
            assert await repository.renew_lease(job_id, "worker-recovery")
            assert await repository.mark_awaiting_configuration(
                job_id,
                "worker-recovery",
                preview={"preview_fingerprint": "a" * 64, "sheets": []},
            )
        async with session_factory() as session:
            repository = KnowledgeImportJobRepository(session)
            profile = await repository.confirm_excel_profile(
                job_id,
                user_id=user_id,
                preview_fingerprint="a" * 64,
                mode="section",
                profile_name="Section fallback",
                sheets=[],
            )
            assert profile.status == "draft"
            assert (await repository.get(job_id)).status == "pending"
    finally:
        async with session_factory() as session:
            await session.execute(
                delete(KnowledgeIngestionJobTable).where(
                    KnowledgeIngestionJobTable.id == job_id
                )
            )
            await session.execute(
                delete(KnowledgeDocumentTable).where(
                    KnowledgeDocumentTable.doc_id
                    == build_doc_id(
                        f"knowledge-base/product_planning/{job_id}.xlsx"
                    )
                )
            )
            await session.commit()
        await engine.dispose()


async def main() -> None:
    """运行无外部依赖检查，并按命令行开关追加真实数据库验收。"""

    with TemporaryDirectory() as directory:
        root = Path(directory)
        test_office_builders_and_incremental_diff()
        test_mixed_excel_profile_and_coordinates()
        test_default_loader_excludes_office()
        test_office_loaders_and_stable_chunks(root)
        test_notes_none_guard(root)
        test_ooxml_rejections(root)
        test_filename_and_request_limits()
        test_api_contract(root)
        await test_replace_docs_idempotency(root)
        await test_periodic_heartbeat()
        test_worker_department_acl(root)
    if "--real-db" in sys.argv:
        await test_postgres_claim_and_recovery()
    if "--real-stores" in sys.argv:
        await test_real_incremental_stores()
    if "--real-worker" in sys.argv:
        await test_real_worker_incremental_update()
    print("office ingestion checks passed")


if __name__ == "__main__":
    asyncio.run(main())
